import argparse
import sys
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from tools.generate_slides import parse_slides


def load_notes(outline_path: Path) -> dict:
    """Load notes by slide number from a YAML outline using parse_slides."""
    notes = {}
    try:
        slides = parse_slides(str(outline_path))
        for slide in slides:
            num = slide.get("number")
            if num is not None:
                note = slide.get("notes") or ""
                if note:
                    notes[num] = note
    except Exception as e:
        print(f"Warning: failed to parse notes from {outline_path}: {e}", file=sys.stderr)
    return notes


def collect_slide_files(run_dir: Path, use_4k: bool) -> dict:
    """Collect best image for each slide number under run_dir."""
    pattern = re.compile(r"slide_(\d+)_0(.*)\.(jpg|jpeg|png)$", re.IGNORECASE)
    slide_candidates = {}

    for path in run_dir.glob("slide_*_0.*"):
        m = pattern.match(path.name)
        if not m:
            continue
        num = int(m.group(1))
        suffix = m.group(2).lower()
        slide_candidates.setdefault(num, []).append((suffix, path))

    slide_files = {}
    for num, items in slide_candidates.items():
        # prefer 4k if requested
        if use_4k:
            fourk = [p for s, p in items if "4k" in s]
            if fourk:
                slide_files[num] = fourk[0]
                continue
        # otherwise prefer base
        base = [p for s, p in items if not s]
        if base:
            slide_files[num] = base[0]
        else:
            # fallback to first candidate
            slide_files[num] = items[0][1]

    return slide_files


def export_pptx(run_dir: Path, output_path: Path, notes_by_number: dict, use_4k: bool):
    slide_files = collect_slide_files(run_dir, use_4k)
    if not slide_files:
        print(f"Error: no slide images found in {run_dir}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    # Set 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[0]

    for num in sorted(slide_files.keys()):
        img_path = slide_files[num]
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if notes_by_number:
            note = notes_by_number.get(num)
            if note:
                notes_slide = slide.notes_slide
                tf = notes_slide.notes_text_frame
                tf.text = note

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Saved PPTX: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="Export generated slides to PPTX (image-only).")
    parser.add_argument("--run-dir", required=True, help="Path to a generation run directory under generated_slides/")
    parser.add_argument("--output", help="Output PPTX path (default: pptx/<run-dir-name>.pptx)")
    parser.add_argument("--yaml", help="Path to YAML outline for speaker notes (optional)")
    parser.add_argument("--use-4k", action="store_true", help="Prefer 4K images if available")
    args = parser.parse_args()

    script_dir = Path(__file__).parent
    project_root = script_dir.parent

    run_dir = Path(args.run_dir)
    if not run_dir.is_absolute():
        run_dir = (project_root / run_dir).resolve()
    if not run_dir.exists() or not run_dir.is_dir():
        print(f"Error: run directory not found: {run_dir}", file=sys.stderr)
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
        if not output_path.is_absolute():
            output_path = (project_root / output_path).resolve()
    else:
        pptx_dir = project_root / "pptx"
        pptx_dir.mkdir(parents=True, exist_ok=True)
        output_path = pptx_dir / f"{run_dir.name}.pptx"

    notes_by_number = {}
    if args.yaml:
        outline_path = Path(args.yaml)
        if not outline_path.is_absolute():
            outline_path = (project_root / outline_path).resolve()
        if outline_path.exists():
            notes_by_number = load_notes(outline_path)

    export_pptx(run_dir, output_path, notes_by_number, args.use_4k)


if __name__ == "__main__":
    main()
